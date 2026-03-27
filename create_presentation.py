#!/usr/bin/env python3
"""Generate the master's dissertation defense presentation (.pptx).

All facts and numbers are taken directly from the LaTeX report chapters.
Modern visual design with equation boxes, stat callouts, card layouts.
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent
FIGURES = ROOT / "report" / "figures"
LOGO_PATH = ROOT / "report" / "uacs.png"
OUTPUT = ROOT / "defense_presentation.pptx"

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
# Primary palette
UA_RED = RGBColor(0xAB, 0x00, 0x33)
DARK = RGBColor(0x1C, 0x1C, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xF8, 0xFA)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xF2)
MED_GRAY = RGBColor(0x88, 0x88, 0x99)
BODY_COLOR = RGBColor(0x2D, 0x2D, 0x3D)

# Accent palette
ACCENT_BLUE = RGBColor(0x2D, 0x6A, 0x9F)
ACCENT_TEAL = RGBColor(0x0D, 0x8A, 0x72)
ACCENT_AMBER = RGBColor(0xCF, 0x7C, 0x00)
ACCENT_RED_LIGHT = RGBColor(0xF4, 0xE6, 0xEB)  # light pink for highlight bg
EQ_BG = RGBColor(0xE8, 0xEC, 0xF4)  # light blue-gray for equation boxes
CARD_BG = RGBColor(0xF2, 0xF3, 0xF7)
CARD_BORDER = RGBColor(0xD0, 0xD3, 0xDE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Calibri Light"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"
FONT_HEADING = "Calibri"

# Logo (1032x292, AR=3.534)
LOGO_H = Inches(0.38)
LOGO_W = Inches(0.38 * 3.534)
LOGO_LEFT = SLIDE_W - LOGO_W - Inches(0.35)
LOGO_TOP = Inches(0.10)

TITLE_BAR_H = Inches(0.85)
CONTENT_TOP = Inches(1.15)
CONTENT_LEFT = Inches(0.8)
CONTENT_W = Inches(11.7)
CONTENT_H = Inches(5.9)

FIG_TOP = Inches(1.25)
FIG_BOTTOM = SLIDE_H - Inches(0.95)
FIG_MAX_H = FIG_BOTTOM - FIG_TOP
FIG_MAX_W = SLIDE_W - Inches(1.6)

# ---------------------------------------------------------------------------
# Presentation object
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ---------------------------------------------------------------------------
# Core helpers
# ---------------------------------------------------------------------------

def _get_image_size(path):
    with Image.open(path) as img:
        return img.size


def _fit_image(fig_path, max_w, max_h):
    w_px, h_px = _get_image_size(fig_path)
    ar = w_px / h_px
    w = max_w
    h = int(w / ar)
    if h > max_h:
        h = max_h
        w = int(h * ar)
    return w, h


def _add_logo(slide):
    slide.shapes.add_picture(str(LOGO_PATH), LOGO_LEFT, LOGO_TOP, LOGO_W, LOGO_H)


def _no_line(shape):
    shape.line.fill.background()


def _rounded_rect(slide, left, top, w, h, fill_rgb, border_rgb=None, radius=Inches(0.08)):
    """Add a rounded rectangle with optional border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(1)
    else:
        _no_line(shape)
    # Adjust corner radius via XML
    try:
        from pptx.oxml.ns import qn
        sp = shape._element
        for prstGeom in sp.iter(qn('a:prstGeom')):
            for avLst in prstGeom.iter(qn('a:avLst')):
                for gd in list(avLst):
                    avLst.remove(gd)
                from lxml import etree
                gd = etree.SubElement(avLst, qn('a:gd'))
                gd.set('name', 'adj')
                gd.set('fmla', 'val 8000')
    except Exception:
        pass
    return shape


def _add_title_bar(slide, text):
    """Dark gradient-style title bar with thin accent underline."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    _no_line(bar)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, TITLE_BAR_H,
                                     SLIDE_W, Inches(0.035))
    accent.fill.solid()
    accent.fill.fore_color.rgb = UA_RED
    _no_line(accent)

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_right = Inches(2.0)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE


def _slide_number_placeholder(slide, num):
    """Subtle slide number in bottom-right."""
    box = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.42),
                                    Inches(0.7), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT


# Track slide count for numbering
_slide_num = [0]


def _new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_num[0] += 1
    return slide


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_section_divider(title, subtitle=""):
    slide = _new_slide()
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.6),
                                  Inches(0.06), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = UA_RED
    _no_line(bar)

    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(11), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.1), Inches(3.8), Inches(11), Inches(0.7))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(19)
        p2.font.color.rgb = MED_GRAY
        p2.font.name = FONT_BODY

    _add_logo(slide)
    _slide_number_placeholder(slide, _slide_num[0])
    return slide


def _content_slide(title, bullets, bullet_font_size=Pt(20)):
    slide = _new_slide()
    _add_title_bar(slide, title)

    txBox = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_W, CONTENT_H)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        is_sub = bullet.startswith("  ")
        text = bullet.strip()
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = FONT_BODY
        p.font.color.rgb = BODY_COLOR

        if is_sub:
            p.text = text
            p.level = 1
            p.font.size = Pt(17)
            p.space_after = Pt(3)
            p.space_before = Pt(1)
        else:
            p.text = text
            p.level = 0
            p.font.size = bullet_font_size
            p.space_after = Pt(7)
            p.space_before = Pt(4)

    _add_logo(slide)
    _slide_number_placeholder(slide, _slide_num[0])
    return slide


def _figure_slide(title, fig_path, caption=""):
    slide = _new_slide()
    _add_title_bar(slide, title)

    if Path(fig_path).exists():
        w, h = _fit_image(fig_path, FIG_MAX_W, FIG_MAX_H)
        left = int((SLIDE_W - w) / 2)
        top = int(FIG_TOP + (FIG_MAX_H - h) / 2)
        slide.shapes.add_picture(str(fig_path), left, top, w, h)

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(1.2), SLIDE_H - Inches(0.7),
                                            SLIDE_W - Inches(2.4), Inches(0.5))
        cap_tf = cap_box.text_frame
        cap_tf.word_wrap = True
        cap_p = cap_tf.paragraphs[0]
        cap_p.text = caption
        cap_p.font.size = Pt(13)
        cap_p.font.italic = True
        cap_p.font.color.rgb = MED_GRAY
        cap_p.font.name = FONT_BODY
        cap_p.alignment = PP_ALIGN.CENTER

    _add_logo(slide)
    _slide_number_placeholder(slide, _slide_num[0])
    return slide


def _table_slide(title, headers, rows, col_widths=None):
    slide = _new_slide()
    _add_title_bar(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(12.0)
    tbl_left = int((SLIDE_W - tbl_width) / 2)
    tbl_top = Inches(1.35)
    row_height = Inches(0.42)
    tbl_height = row_height * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        cw = int(tbl_width / n_cols)
        for i in range(n_cols):
            table.columns[i].width = cw

    def _fmt(cell, text, bold=False, bg=None, fg=None, font_size=Pt(14)):
        cell.text = str(text)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.name = FONT_BODY
            p.font.bold = bold
            p.font.color.rgb = fg or (WHITE if bg == DARK else BODY_COLOR)
            p.alignment = PP_ALIGN.CENTER
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    for j, h in enumerate(headers):
        _fmt(table.cell(0, j), h, bold=True, bg=DARK, fg=WHITE, font_size=Pt(14))

    for i, row in enumerate(rows):
        bg = LIGHT_GRAY if i % 2 == 0 else None
        for j, val in enumerate(row):
            _fmt(table.cell(i + 1, j), val, bg=bg)

    _add_logo(slide)
    _slide_number_placeholder(slide, _slide_num[0])
    return slide


def _two_figure_slide(title, fig_left_path, fig_right_path, cap_left="", cap_right=""):
    slide = _new_slide()
    _add_title_bar(slide, title)

    half_w = Inches(5.8)
    max_h = FIG_MAX_H
    gap = Inches(0.3)

    for fig_path, region_left, cap in [
        (fig_left_path, Inches(0.5), cap_left),
        (fig_right_path, Inches(0.5) + half_w + gap, cap_right),
    ]:
        if Path(fig_path).exists():
            w, h = _fit_image(fig_path, half_w, max_h)
            left = int(region_left + (half_w - w) / 2)
            top = int(FIG_TOP + (max_h - h) / 2)
            slide.shapes.add_picture(str(fig_path), left, top, w, h)
        if cap:
            cb = slide.shapes.add_textbox(region_left, SLIDE_H - Inches(0.7),
                                           half_w, Inches(0.5))
            ctf = cb.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.text = cap
            cp.font.size = Pt(12)
            cp.font.italic = True
            cp.font.color.rgb = MED_GRAY
            cp.alignment = PP_ALIGN.CENTER

    _add_logo(slide)
    _slide_number_placeholder(slide, _slide_num[0])
    return slide


# ---------------------------------------------------------------------------
# Modern visual helpers
# ---------------------------------------------------------------------------

def _equation_box(slide, left, top, width, height, text, font_size=Pt(18)):
    """A rounded box with monospace equation text on a light background."""
    _rounded_rect(slide, left, top, width, height, EQ_BG, CARD_BORDER)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08),
                                    width - Inches(0.4), height - Inches(0.16))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = FONT_MONO
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    return box


def _highlight_box(slide, left, top, width, height, text, bg=ACCENT_RED_LIGHT,
                   font_size=Pt(16), font_color=None, bold=False):
    """A tinted highlight box for callouts and key points."""
    _rounded_rect(slide, left, top, width, height, bg)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.06),
                                    width - Inches(0.4), height - Inches(0.12))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = FONT_BODY
    p.font.color.rgb = font_color or BODY_COLOR
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT
    return tf


def _stat_card(slide, left, top, width, height, number, label,
               num_color=UA_RED, num_size=Pt(44)):
    """A stat callout card: big number + description underneath."""
    _rounded_rect(slide, left, top, width, height, CARD_BG, CARD_BORDER)
    # Number
    nb = slide.shapes.add_textbox(left, top + Inches(0.15), width, Inches(0.8))
    ntf = nb.text_frame
    np_ = ntf.paragraphs[0]
    np_.text = str(number)
    np_.font.size = num_size
    np_.font.bold = True
    np_.font.color.rgb = num_color
    np_.font.name = FONT_HEADING
    np_.alignment = PP_ALIGN.CENTER
    # Label
    lb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.95),
                                   width - Inches(0.3), height - Inches(1.1))
    ltf = lb.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(13)
    lp.font.color.rgb = MED_GRAY
    lp.font.name = FONT_BODY
    lp.alignment = PP_ALIGN.CENTER


def _card_with_title(slide, left, top, width, height, card_title, bullets,
                     title_color=ACCENT_BLUE):
    """A framed card with a colored title and bullet list."""
    _rounded_rect(slide, left, top, width, height, WHITE, CARD_BORDER)
    # Colored top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = title_color
    _no_line(bar)

    # Title
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                   width - Inches(0.4), Inches(0.35))
    ttf = tb.text_frame
    tp = ttf.paragraphs[0]
    tp.text = card_title
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = title_color
    tp.font.name = FONT_HEADING

    # Bullets
    bb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.5),
                                   width - Inches(0.5), height - Inches(0.65))
    btf = bb.text_frame
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_COLOR
        p.font.name = FONT_BODY
        p.space_after = Pt(4)


# ===================================================================
#  SLIDES
# ===================================================================

# ------ 1. TITLE SLIDE ------
slide = _new_slide()
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK

# Subtle gradient-like stripe
stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0), Inches(0.12), SLIDE_H)
stripe.fill.solid()
stripe.fill.fore_color.rgb = UA_RED
_no_line(stripe)

# Accent line
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.35), SLIDE_W, Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = UA_RED
_no_line(bar)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.0), Inches(2.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Adapting SimLingo Vision-Language-Action"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT_TITLE

p1b = tf.add_paragraph()
p1b.text = "Models to New Domains"
p1b.font.size = Pt(40)
p1b.font.bold = True
p1b.font.color.rgb = WHITE
p1b.font.name = FONT_TITLE
p1b.space_before = Pt(2)

p2 = tf.add_paragraph()
p2.text = "A Case Study for QCar 2"
p2.font.size = Pt(24)
p2.font.color.rgb = UA_RED
p2.font.name = FONT_HEADING
p2.font.bold = True
p2.space_before = Pt(20)

info = slide.shapes.add_textbox(Inches(1.0), Inches(4.65), Inches(11), Inches(2.4))
itf = info.text_frame
itf.word_wrap = True
for line, sz, color in [
    ("Garegin Mazmanyan", Pt(24), WHITE),
    ("Master of Science in Computer Science", Pt(17), MED_GRAY),
    ("The University of Arizona", Pt(17), MED_GRAY),
    ("", Pt(8), WHITE),
    ("Committee:  Chicheng Zhang (Chair)  \u00b7  Hossein Rastgoftar (Co-Chair)  \u00b7  Eduardo Blanco", Pt(14), MED_GRAY),
    ("February 2026", Pt(14), MED_GRAY),
]:
    pi = itf.add_paragraph()
    pi.text = line
    pi.font.size = sz
    pi.font.color.rgb = color
    pi.font.name = FONT_BODY
    pi.space_after = Pt(3)

_add_logo(slide)


# ------ 2. OUTLINE ------
slide = _new_slide()
_add_title_bar(slide, "Outline")

sections = [
    ("01", "Motivation & Research Questions", ACCENT_BLUE),
    ("02", "Background: SimLingo & VLA Driving", ACCENT_BLUE),
    ("03", "Methodology: CARLA \u2192 QLabs", ACCENT_TEAL),
    ("04", "Data Collection & Fine-Tuning", ACCENT_TEAL),
    ("05", "Inference Stack & Control", ACCENT_TEAL),
    ("06", "Evaluation Design & Baselines", ACCENT_AMBER),
    ("07", "Results & Analysis", UA_RED),
    ("08", "Limitations & Future Work", MED_GRAY),
    ("09", "Conclusion", DARK),
]
y = Inches(1.25)
for num, label, color in sections:
    # Number badge
    _rounded_rect(slide, Inches(1.0), y, Inches(0.65), Inches(0.45), color)
    nb = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.02), Inches(0.65), Inches(0.4))
    ntf = nb.text_frame
    np_ = ntf.paragraphs[0]
    np_.text = num
    np_.font.size = Pt(16)
    np_.font.bold = True
    np_.font.color.rgb = WHITE
    np_.font.name = FONT_HEADING
    np_.alignment = PP_ALIGN.CENTER
    # Label
    lb = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.03), Inches(9), Inches(0.4))
    ltf = lb.text_frame
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(20)
    lp.font.color.rgb = BODY_COLOR
    lp.font.name = FONT_BODY
    y += Inches(0.6)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 3. MOTIVATION ------
slide = _new_slide()
_add_title_bar(slide, "Motivation")

# Left column: text
bullets_left = [
    "Vision-Language-Action (VLA) models show\npromise for autonomous driving",
    "SimLingo: winning entry of the CARLA\nAutonomous Driving Challenge 2024",
    "Uses InternVL2-1B multimodal backbone\nfor end-to-end waypoint prediction",
]
y = Inches(1.3)
for b in bullets_left:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(6.5), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(18)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    p.space_after = Pt(4)
    y += Inches(0.95)

# Right column: key question card
_rounded_rect(slide, Inches(7.8), Inches(1.3), Inches(4.8), Inches(2.8), CARD_BG, CARD_BORDER)
qb = slide.shapes.add_textbox(Inches(8.1), Inches(1.5), Inches(4.3), Inches(0.5))
qtf = qb.text_frame
qp = qtf.paragraphs[0]
qp.text = "Core Question"
qp.font.size = Pt(16)
qp.font.bold = True
qp.font.color.rgb = UA_RED
qp.font.name = FONT_HEADING

qb2 = slide.shapes.add_textbox(Inches(8.1), Inches(2.1), Inches(4.3), Inches(1.8))
qtf2 = qb2.text_frame
qtf2.word_wrap = True
qp2 = qtf2.paragraphs[0]
qp2.text = "Can a VLA driving model trained on CARLA transfer to Quanser QLabs via parameter-efficient fine-tuning?"
qp2.font.size = Pt(17)
qp2.font.color.rgb = BODY_COLOR
qp2.font.name = FONT_BODY

# Bottom highlight
_highlight_box(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.65),
               "Goal: reproducible pipeline for adapting CARLA-trained VLA to QCar 2 in QLabs",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(18), font_color=ACCENT_BLUE, bold=True)

# Comparison cards
_card_with_title(slide, Inches(0.8), Inches(5.5), Inches(5.6), Inches(1.6),
                 "CARLA", ["Unreal Engine 4, photorealistic rendering",
                           "Full vehicle dynamics (tire, suspension)",
                           "Rich sensor suite"], title_color=ACCENT_BLUE)
_card_with_title(slide, Inches(6.9), Inches(5.5), Inches(5.6), Inches(1.6),
                 "QLabs", ["Educational platform, simplified rendering",
                           "Idealized physics (instant velocity)",
                           "Compact fixed-layout environment"], title_color=ACCENT_TEAL)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ SECTION: BACKGROUND ------
_add_section_divider("Background", "SimLingo, InternVL2-1B, and VLA Driving Architecture")


# ------ 4. TWO-LEVEL ARCHITECTURE ------
slide = _new_slide()
_add_title_bar(slide, "Two-Level Driving Architecture")

_card_with_title(slide, Inches(0.8), Inches(1.25), Inches(5.6), Inches(2.7),
                 "Upper Level (Learned)", [
                     "VLA model predicts future waypoints",
                     "Input: camera image + route intent",
                     "Output: WHERE the vehicle should go",
                     "Never directly outputs steering or throttle",
                 ], title_color=ACCENT_BLUE)

_card_with_title(slide, Inches(6.9), Inches(1.25), Inches(5.6), Inches(2.7),
                 "Lower Level (Classical)", [
                     "PID controllers convert waypoints to actuators",
                     "Lateral PID: heading error \u2192 steering angle",
                     "Longitudinal PID: speed error \u2192 throttle",
                     "Robust, interpretable control layer",
                 ], title_color=ACCENT_TEAL)

# Arrow between cards
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(2.3),
                                Inches(0.9), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = MED_GRAY
_no_line(arrow)

_highlight_box(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.55),
               "Clean separation: learned perception + classical control",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(18), font_color=ACCENT_BLUE, bold=True)

# Architecture flow
_content_slide.__wrapped__ = True  # dummy - we're adding extra content manually
steps = ["Camera\nImage", "Vision\nEncoder", "LLM\nDecoder", "Waypoint\nHeads", "PID\nControllers", "Vehicle\nCommands"]
colors = [ACCENT_BLUE, ACCENT_BLUE, ACCENT_BLUE, ACCENT_TEAL, ACCENT_TEAL, ACCENT_TEAL]
x_start = Inches(0.8)
box_w = Inches(1.6)
gap = Inches(0.35)
y_flow = Inches(5.2)
for idx, (step, color) in enumerate(zip(steps, colors)):
    x = int(x_start + idx * (box_w + gap))
    _rounded_rect(slide, x, y_flow, box_w, Inches(1.0), color)
    tb = slide.shapes.add_textbox(x, y_flow + Inches(0.1), box_w, Inches(0.8))
    ttf = tb.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = step
    tp.font.size = Pt(13)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.font.name = FONT_HEADING
    tp.alignment = PP_ALIGN.CENTER
    # Arrow
    if idx < len(steps) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      int(x + box_w), y_flow + Inches(0.3),
                                      gap, Inches(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = LIGHT_GRAY
        _no_line(arr)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 5. SimLingo & InternVL2-1B ------
slide = _new_slide()
_add_title_bar(slide, "SimLingo & InternVL2-1B Backbone")

# Stat cards row
_stat_card(slide, Inches(0.8), Inches(1.2), Inches(2.6), Inches(1.5),
           "300M", "Vision encoder\nInternViT-300M-448px", num_color=ACCENT_BLUE)
_stat_card(slide, Inches(3.7), Inches(1.2), Inches(2.6), Inches(1.5),
           "500M", "Language model\nQwen2-0.5B-Instruct", num_color=ACCENT_TEAL)
_stat_card(slide, Inches(6.6), Inches(1.2), Inches(2.6), Inches(1.5),
           "512", "Visual tokens\n(2 tiles \u00d7 256 each)", num_color=ACCENT_AMBER)
_stat_card(slide, Inches(9.5), Inches(1.2), Inches(2.6), Inches(1.5),
           "448\u00b2", "Tile resolution\npx per tile", num_color=UA_RED)

# Details below
_content_slide_bullets = [
    "Dynamic tiling: images split into 448\u00d7448 tiles, pixel-unshuffle (factor 4) \u2192 256 tokens/tile",
    "Learned projection maps vision features into LLM embedding space",
    "Causal LLM processes visual + text + navigational tokens jointly",
]
y = Inches(3.1)
for b in _content_slide_bullets:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(18)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    p.space_after = Pt(4)
    y += Inches(0.55)

# Architecture equation
_equation_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.55),
              "Token sequence:  [NAV_1, NAV_2]  +  [VIS_1 ... VIS_512]  +  [TEXT tokens]  \u2192  LLM  \u2192  Waypoint Heads",
              font_size=Pt(15))

_highlight_box(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.55),
               "Target-point mode: 2 GPS-like waypoints fed through WaypointInputAdaptor MLP \u2192 2 navigational embeddings",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(16), font_color=ACCENT_BLUE)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 6. ACTION PREDICTION ------
slide = _new_slide()
_add_title_bar(slide, "Dual-Head Action Prediction")

_card_with_title(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(2.8),
                 "Speed Waypoints (Temporal)", [
                     "10 learnable speed query vectors",
                     "Predict ego-frame coords at 0.25 s intervals",
                     "Used for longitudinal PID control",
                     "MLP: Linear(H\u2192256) \u2192 SiLU \u2192 Linear(256\u21922)",
                 ], title_color=ACCENT_BLUE)

_card_with_title(slide, Inches(6.9), Inches(1.2), Inches(5.6), Inches(2.8),
                 "Route Waypoints (Geometric)", [
                     "20 learnable route query vectors",
                     "Predict ego-frame coords at ~1 m intervals",
                     "Used for lateral PID control",
                     "MLP: Linear(H\u2192512) \u2192 SiLU \u2192 ... \u2192 Linear(256\u21922)",
                 ], title_color=ACCENT_TEAL)

_equation_box(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.55),
              "Output:  waypoint deltas  \u2192  cumulative sum  \u2192  absolute ego-frame positions",
              font_size=Pt(16))

_highlight_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.55),
               "Training losses:  Smooth-L1 (route + speed waypoints) + Cross-entropy (language)  \u2014  equal weight",
               bg=ACCENT_RED_LIGHT, font_size=Pt(16), font_color=BODY_COLOR, bold=True)

# MLP adaptor
_equation_box(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.55),
              "WaypointInputAdaptor:  Linear(2\u2192256) \u2192 ReLU \u2192 Linear(256\u2192512) \u2192 ReLU \u2192 Linear(512\u2192D)",
              font_size=Pt(14))

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 8. CARLA vs QLabs ------
_table_slide(
    "CARLA vs. QLabs: Key Differences",
    ["Category", "CARLA", "QLabs"],
    [
        ["Engine", "Unreal Engine 4, photorealistic", "Educational, simplified rendering"],
        ["Physics", "Full dynamics (tire, suspension)", "Idealized (instant velocity)"],
        ["Steering sign", "Positive = turn left", "Opposite sign convention"],
        ["Control cadence", "20 Hz control, 4 Hz data", "4 Hz control loop"],
        ["Speed signal", "Instantaneous from simulator", "Single-frame displacement"],
        ["Camera", "1024\u00d7512, FOV 110\u00b0", "820\u00d7410, FOV 160\u00b0"],
    ],
    col_widths=[Inches(2.5), Inches(4.75), Inches(4.75)],
)


# ------ 9. RESEARCH QUESTIONS ------
slide = _new_slide()
_add_title_bar(slide, "Research Questions")

rqs = [
    ("RQ1", "What data and interface adaptations are needed for CARLA \u2192 QLabs transfer?",
     "Coordinate frames, steering signs, camera configs, speed estimation", ACCENT_BLUE),
    ("RQ2", "Can parameter-efficient fine-tuning (LoRA) achieve domain transfer?",
     "Freeze vision encoder, adapt only LLM parameters", ACCENT_TEAL),
    ("RQ3", "What runtime alignment is required for real-time control?",
     "Control-loop timing, PID compensation, ego-frame target points", ACCENT_AMBER),
    ("RQ4", "How to build a reproducible evaluation harness?",
     "Scenario-driven testing with quantitative metrics", UA_RED),
]
y = Inches(1.2)
for tag, question, detail, color in rqs:
    # Tag badge
    _rounded_rect(slide, Inches(0.8), y, Inches(0.85), Inches(1.2), color)
    tb = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.3), Inches(0.85), Inches(0.5))
    ttf = tb.text_frame
    tp = ttf.paragraphs[0]
    tp.text = tag
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.font.name = FONT_HEADING
    tp.alignment = PP_ALIGN.CENTER
    # Question
    qbox = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.05), Inches(10.5), Inches(0.55))
    qtf = qbox.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = question
    qp.font.size = Pt(18)
    qp.font.bold = True
    qp.font.color.rgb = BODY_COLOR
    qp.font.name = FONT_BODY
    # Detail
    dbox = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.6), Inches(10.5), Inches(0.45))
    dtf = dbox.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = detail
    dp.font.size = Pt(15)
    dp.font.color.rgb = MED_GRAY
    dp.font.name = FONT_BODY
    y += Inches(1.4)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ SECTION: METHODOLOGY ------
_add_section_divider("Methodology", "Adapting CARLA \u2192 QLabs")


# ------ 10. DOMAIN ADAPTATION OVERVIEW ------
slide = _new_slide()
_add_title_bar(slide, "Domain Adaptation Pipeline")

_highlight_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
               "Starting point: pretrained CARLA checkpoint (epoch 13 of CARLA training)",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(17), font_color=ACCENT_BLUE, bold=True)

# Three pipeline stages as cards
stages = [
    ("Stage 1: Data Collection", ["Expert driving in QLabs", "Keyboard teleop at 30 Hz", "4 Hz data logging"], ACCENT_BLUE),
    ("Stage 2: LoRA Fine-Tuning", ["InternVL2-1B backbone", "~17.6M adapter params", "15 epochs training"], ACCENT_TEAL),
    ("Stage 3: Inference Stack", ["Real-time PID controller", "Ego-frame target points", "4 Hz control loop"], ACCENT_AMBER),
]
for i, (stitle, sbullets, scolor) in enumerate(stages):
    x = Inches(0.8) + i * Inches(4.1)
    _card_with_title(slide, x, Inches(2.0), Inches(3.8), Inches(2.2),
                     stitle, sbullets, title_color=scolor)

# Key mismatches
_content_slide_text = "Key mismatches addressed:  coordinate frames  \u00b7  control rates (20 Hz vs 4 Hz)  \u00b7  camera FOV/resolution  \u00b7  speed estimation"
_highlight_box(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5),
               _content_slide_text, bg=ACCENT_RED_LIGHT, font_size=Pt(15))

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 11. COORDINATE FRAME ------
slide = _new_slide()
_add_title_bar(slide, "Coordinate Frame Alignment")

bullets = [
    "Ego frame: x points forward, y points left",
    "All waypoints converted to ego frame before use",
    "Steering sign: CARLA and QLabs use opposite conventions",
]
y = Inches(1.2)
for b in bullets:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(19)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    y += Inches(0.55)

_equation_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
              "World \u2192 Ego:    p\u2091 = R(\u2212\u03c8) \u00b7 (p\u02b7 \u2212 t)        where R is 2D rotation matrix, \u03c8 = heading")

_equation_box(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
              "Steering sign flip:    \u03b4 = \u2212u \u00b7 \u03b4\u2098\u2090\u2093        (negate for QLabs convention)")

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 12. CONTROL TIMING ------
slide = _new_slide()
_add_title_bar(slide, "Control Timing Compensation")

_card_with_title(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(1.8),
                 "CARLA", ["20 Hz control loop", "Data saved at 4 Hz (every 5th frame)",
                           "PID tuned for 20 Hz feedback"], title_color=ACCENT_BLUE)
_card_with_title(slide, Inches(6.9), Inches(1.2), Inches(5.6), Inches(1.8),
                 "QLabs", ["4 Hz control loop directly", "Derivative term needs compensation",
                           "Without fix: steering oscillation"], title_color=ACCENT_TEAL)

_equation_box(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.65),
              "Derivative compensation:    raw_diff / (0.25s / 0.05s)  =  raw_diff / 5",
              font_size=Pt(17))

_highlight_box(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
               "Normalizes the 4 Hz derivative term to the 20 Hz reference used during CARLA training",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(16), font_color=ACCENT_BLUE)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 13. CAMERA CONFIGURATION ------
_table_slide(
    "Camera Configuration",
    ["Property", "CARLA", "QLabs"],
    [
        ["Mounting position", "[\u22121.5, 0.0, 2.0]", "[0.183, 0.0, 0.110]"],
        ["Resolution", "1024\u00d7512", "820\u00d7410"],
        ["Field of view", "110\u00b0", "160\u00b0"],
        ["Max tiles (InternVL2)", "2", "2"],
        ["Visual tokens", "512", "512"],
    ],
    col_widths=[Inches(3.5), Inches(4.25), Inches(4.25)],
)


# ------ 14. SPEED & TARGET POINTS ------
slide = _new_slide()
_add_title_bar(slide, "Speed Estimation & Target-Point Routing")

_card_with_title(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(2.5),
                 "Speed Estimation", [
                     "Single-frame displacement method",
                     "Same in training & inference \u2192 consistent",
                     "Domain-agnostic (no simulator velocity API)",
                     "Trade-off: noisier but portable",
                 ], title_color=ACCENT_BLUE)

_card_with_title(slide, Inches(6.9), Inches(1.2), Inches(5.6), Inches(2.5),
                 "Target-Point Routing", [
                     "85-waypoint roundabout route (JSON config)",
                     "Advance along polyline until arc-length > d",
                     "Convert target points world \u2192 ego frame",
                     "Feed to WaypointInputAdaptor MLP",
                 ], title_color=ACCENT_TEAL)

_equation_box(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.6),
              "v\u209c = ||p\u209c \u2212 p\u209c\u208b\u2081|| / \u0394t        where \u0394t = 0.25 s  (at 4 Hz)",
              font_size=Pt(17))

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ SECTION: DATA & TRAINING ------
_add_section_divider("Data Collection & Training", "Dataset, LoRA Fine-Tuning, and Convergence")


# ------ 16. DATA COLLECTION ------
_content_slide("Data Collection Pipeline", [
    "Expert driving via keyboard teleoperation at 30 Hz",
    "Data logging at 4 Hz (\u0394t = 0.25 s), matching runtime controller frequency",
    "  Each frame: 1 RGB JPEG + 1 gzip JSON measurement (position, heading, speed)",
    "Route: roundabout navigation with optional single static obstacle",
    "Multi-actor challenge: QLabs API calls are serialized",
    "  Shared mutex prevents simultaneous API calls",
    "  Actor update rate lowered to 2 Hz, start times staggered 0\u201399 ms",
    "  Container timeout extended to 10 s",
])


# ------ 18. DATASET SUMMARY (stat cards + table) ------
slide = _new_slide()
_add_title_bar(slide, "Dataset Summary")

_stat_card(slide, Inches(0.8), Inches(1.2), Inches(2.3), Inches(1.5),
           "10,495", "Total frames\nacross 55 runs", num_color=ACCENT_BLUE, num_size=Pt(36))
_stat_card(slide, Inches(3.4), Inches(1.2), Inches(2.3), Inches(1.5),
           "40 / 15", "Train / Val runs", num_color=ACCENT_TEAL, num_size=Pt(36))
_stat_card(slide, Inches(6.0), Inches(1.2), Inches(2.3), Inches(1.5),
           "~470 MiB", "On-disk footprint", num_color=ACCENT_AMBER, num_size=Pt(34))
_stat_card(slide, Inches(8.6), Inches(1.2), Inches(2.3), Inches(1.5),
           "5", "Obstacle placement\npositions", num_color=UA_RED, num_size=Pt(40))

# Mini table
headers = ["Split", "Runs", "Frames", "Run Length Range"]
rows = [
    ["Training", "40", "7,498", "98 \u2013 301 frames"],
    ["Validation", "15", "2,997", "171 \u2013 234 frames"],
]
n_rows = len(rows) + 1
n_cols = len(headers)
tbl_w = Inches(10.0)
tbl_l = int((SLIDE_W - tbl_w) / 2)
shape = slide.shapes.add_table(n_rows, n_cols, tbl_l, Inches(3.2), tbl_w, Inches(1.2))
table = shape.table
for ci in range(n_cols):
    table.columns[ci].width = int(tbl_w / n_cols)

for j, h in enumerate(headers):
    c = table.cell(0, j)
    c.text = h
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    c.fill.solid()
    c.fill.fore_color.rgb = DARK
    for pp in c.text_frame.paragraphs:
        pp.font.size = Pt(14)
        pp.font.bold = True
        pp.font.color.rgb = WHITE
        pp.font.name = FONT_BODY
        pp.alignment = PP_ALIGN.CENTER

for i, row in enumerate(rows):
    for j, val in enumerate(row):
        c = table.cell(i + 1, j)
        c.text = val
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        if i % 2 == 0:
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT_GRAY
        for pp in c.text_frame.paragraphs:
            pp.font.size = Pt(14)
            pp.font.name = FONT_BODY
            pp.font.color.rgb = BODY_COLOR
            pp.alignment = PP_ALIGN.CENTER

_highlight_box(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.45),
               "Each frame: 1 RGB JPEG + 1 gzip-compressed JSON measurement  \u00b7  Obstacle waypoint indices: 25, 35, 50, 65, 75",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(14), font_color=ACCENT_BLUE)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 19. QLABS ENVIRONMENT (MAP) ------
_figure_slide(
    "QLabs Roundabout Environment",
    FIGURES / "SDCS_RoadMap_RightHandTraffic.png",
    "SDCS road map with numbered waypoint nodes. Roundabout route used for training and evaluation.",
)


# ------ 20. LORA FINE-TUNING ------
slide = _new_slide()
_add_title_bar(slide, "LoRA Fine-Tuning Strategy")

_stat_card(slide, Inches(0.8), Inches(1.2), Inches(2.6), Inches(1.5),
           "17.6M", "LoRA adapter\nparameters", num_color=ACCENT_BLUE, num_size=Pt(36))
_stat_card(slide, Inches(3.7), Inches(1.2), Inches(2.6), Inches(1.5),
           "r = 32", "LoRA rank\n\u03b1 = 64, dropout = 0.1", num_color=ACCENT_TEAL, num_size=Pt(36))
_stat_card(slide, Inches(6.6), Inches(1.2), Inches(2.6), Inches(1.5),
           "Frozen", "300M-param vision\nencoder (InternViT)", num_color=ACCENT_AMBER, num_size=Pt(34))
_stat_card(slide, Inches(9.5), Inches(1.2), Inches(2.6), Inches(1.5),
           "Epoch 13", "CARLA pretrained\nstarting checkpoint", num_color=UA_RED, num_size=Pt(30))

bullets = [
    "Applied to all linear layers of Qwen2-0.5B language model",
    "No language/commentary supervision: constant command = 4 (\u201cfollow the road\u201d)",
    "Trained exclusively in target-point navigation mode",
]
y = Inches(3.1)
for b in bullets:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(18)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    y += Inches(0.5)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 21. TRAINING SETUP ------
_table_slide(
    "Training Hyperparameters",
    ["Parameter", "Value"],
    [
        ["Backbone", "InternVL2-1B (OpenGVLab)"],
        ["Learning rate", "2e-5"],
        ["Optimizer", "AdamW (\u03b2 = [0.9, 0.999], wd = 0.1)"],
        ["Batch size", "4 (effective: 16 with grad accum 4)"],
        ["Max epochs", "15"],
        ["Precision", "16-mixed"],
        ["Strategy", "DeepSpeed Stage 2"],
        ["GPU", "NVIDIA RTX 5070 Ti (16 GB VRAM)"],
        ["LR schedule", "One-cycle with 5% warmup"],
        ["Seed", "42"],
    ],
    col_widths=[Inches(5.0), Inches(7.0)],
)


# ------ 22. TRAINING CONVERGENCE ------
_figure_slide(
    "Training Loss Curves",
    FIGURES / "metrics.png",
    "Total, route, speed, and language loss over 15 epochs. Rapid convergence in first 3\u20135 epochs.",
)


# ------ 22b. ADE CURVE ------
_figure_slide(
    "Average Displacement Error (ADE) Convergence",
    FIGURES / "policy_vs_expert_curve.png",
    "25% ADE reduction: 0.114 m (epoch 0) \u2192 0.085 m (epoch 11). Expert baseline ADE = 0.087 m.",
)


# ------ SECTION: INFERENCE ------
_add_section_divider("Inference Stack", "Runtime Pipeline and Control Conversion")


# ------ 23. RUNTIME PIPELINE ------
slide = _new_slide()
_add_title_bar(slide, "Runtime Inference Pipeline")

steps_data = [
    ("1", "Capture", "Camera frame\nfrom CSI", ACCENT_BLUE),
    ("2", "Ego State", "Position, heading\nspeed estimate", ACCENT_BLUE),
    ("3", "Target Pts", "Ego-frame from\nactive route", ACCENT_TEAL),
    ("4", "VLA Infer", "Or reuse cached\nprediction", ACCENT_TEAL),
    ("5", "Commands", "Waypoints \u2192\nQCar 2 cmds", ACCENT_AMBER),
]
x = Inches(0.6)
box_w = Inches(2.2)
gap = Inches(0.25)
for i, (num, label, desc, color) in enumerate(steps_data):
    cx = int(x + i * (box_w + gap))
    _rounded_rect(slide, cx, Inches(1.3), box_w, Inches(2.4), CARD_BG, CARD_BORDER)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.75), Inches(1.45),
                                   Inches(0.65), Inches(0.65))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    _no_line(circ)
    nb = slide.shapes.add_textbox(cx + Inches(0.75), Inches(1.5), Inches(0.65), Inches(0.55))
    ntf = nb.text_frame
    np_ = ntf.paragraphs[0]
    np_.text = num
    np_.font.size = Pt(22)
    np_.font.bold = True
    np_.font.color.rgb = WHITE
    np_.font.name = FONT_HEADING
    np_.alignment = PP_ALIGN.CENTER
    # Label
    lb = slide.shapes.add_textbox(cx + Inches(0.1), Inches(2.25), box_w - Inches(0.2), Inches(0.4))
    ltf = lb.text_frame
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = BODY_COLOR
    lp.font.name = FONT_HEADING
    lp.alignment = PP_ALIGN.CENTER
    # Desc
    db = slide.shapes.add_textbox(cx + Inches(0.1), Inches(2.7), box_w - Inches(0.2), Inches(0.8))
    dtf = db.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(13)
    dp.font.color.rgb = MED_GRAY
    dp.font.name = FONT_BODY
    dp.alignment = PP_ALIGN.CENTER
    # Arrow
    if i < len(steps_data) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      int(cx + box_w), Inches(2.3), gap, Inches(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = LIGHT_GRAY
        _no_line(arr)

_highlight_box(slide, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.5),
               "Nominal control rate: 4 Hz (\u0394t = 0.25 s)  \u00b7  Inference stride configurable  \u00b7  All eval runs: stride = 1",
               bg=RGBColor(0xE8, 0xF0, 0xF8), font_size=Pt(16), font_color=ACCENT_BLUE)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 24. CONTROL CONVERSION ------
slide = _new_slide()
_add_title_bar(slide, "Control Conversion: Speed & Steering")

_card_with_title(slide, Inches(0.8), Inches(1.2), Inches(5.6), Inches(2.0),
                 "Longitudinal (Speed)", [
                     "Distance between first two speed waypoints / \u0394t",
                     "Clipped to [0, v_max] for safety",
                     "Sent directly to QCar 2 velocity interface",
                 ], title_color=ACCENT_BLUE)

_card_with_title(slide, Inches(6.9), Inches(1.2), Inches(5.6), Inches(2.0),
                 "Lateral (Steering PID)", [
                     "Heading error scaled: rad \u2192 deg \u2192 /90",
                     "Derivative: raw_diff / 5 (4 Hz \u2192 20 Hz ref)",
                     "Integral: K_I = 0 (currently inactive)",
                 ], title_color=ACCENT_TEAL)

_equation_box(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.55),
              "v_des = clip( ||s\u2081 \u2212 s\u2080|| / \u0394t_model ,  0,  v_max )",
              font_size=Pt(17))

_equation_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.55),
              "\u03b4 = \u2212u \u00b7 \u03b4_max        (steering sign flip for QLabs convention)",
              font_size=Pt(17))

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ SECTION: EVALUATION ------
_add_section_divider("Evaluation", "Test Design, Metrics, and Baselines")


# ------ 27. EVALUATION DESIGN ------
slide = _new_slide()
_add_title_bar(slide, "Evaluation Design")

_stat_card(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(1.5),
           "15", "Total test runs\nper full sweep", num_color=ACCENT_BLUE)
_stat_card(slide, Inches(4.7), Inches(1.2), Inches(3.5), Inches(1.5),
           "5 + 10", "Baseline + obstacle\nruns", num_color=ACCENT_TEAL)
_stat_card(slide, Inches(8.6), Inches(1.2), Inches(3.5), Inches(1.5),
           "5", "Obstacle placement\nvariants", num_color=ACCENT_AMBER)

_content_slide_bullets2 = [
    "Baseline (no obstacle): 5 repeated runs",
    "Static obstacle variants 1\u20135: 2 repeated runs each",
    "Placements span: early roundabout, mid, exit, straight, late (curved)",
]
y = Inches(3.0)
for b in _content_slide_bullets2:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(18)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    y += Inches(0.5)

_highlight_box(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.55),
               "Caveat: test obstacle positions are identical to training positions \u2014 this is an in-distribution reliability test, not a generalization test",
               bg=ACCENT_RED_LIGHT, font_size=Pt(15), font_color=UA_RED, bold=True)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 28. METRICS ------
slide = _new_slide()
_add_title_bar(slide, "Evaluation Metrics")

metrics = [
    ("Lateral Deviation", "Cross-track error: point-to-polyline-segment distance\nReports both mean and max deviation", ACCENT_BLUE),
    ("Route Coverage", "Waypoint reached within r = 1.5 m\nCoverage = reached / total \u00d7 100%", ACCENT_TEAL),
    ("Safety", "Bumper-hit flags from QLabs API\nStop detection at speed < 0.05 m/s", ACCENT_AMBER),
    ("Timeout", "Max step count + forward-progress thresholds\nPrevents infinite loops", MED_GRAY),
]
for i, (mtitle, mdesc, mcolor) in enumerate(metrics):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.2) + row * Inches(2.5)
    _card_with_title(slide, x, y, Inches(5.8), Inches(2.1), mtitle,
                     mdesc.split("\n"), title_color=mcolor)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 29. ACC BASELINE ------
_figure_slide(
    "ACC Baseline: LiDAR Filtering Pipeline",
    FIGURES / "lidar_filtering_explained.png",
    "4-step LiDAR filter: raw \u2192 forward cone \u2192 world frame \u2192 lane boundary. Pure-pursuit lateral with 4.5 m lookahead.",
)


# ------ 30. OBSTACLE MAP ------
_figure_slide(
    "Obstacle Placement Map",
    FIGURES / "roundabout_obstacle_map.png",
    "5 obstacle variants: green = clean stop, orange = low-speed contact, red = collision.",
)


# ------ SECTION: RESULTS ------
_add_section_divider("Results", "Training Convergence, Route Following, and Obstacle Avoidance")


# ------ 31. TRAINING RESULTS ------
slide = _new_slide()
_add_title_bar(slide, "Training Results")

_stat_card(slide, Inches(0.8), Inches(1.2), Inches(3.0), Inches(1.7),
           "25%", "ADE reduction\nover 15 epochs", num_color=ACCENT_BLUE, num_size=Pt(48))
_stat_card(slide, Inches(4.1), Inches(1.2), Inches(3.0), Inches(1.7),
           "0.085 m", "Final ADE\n(< 9 cm avg deviation)", num_color=ACCENT_TEAL, num_size=Pt(36))
_stat_card(slide, Inches(7.4), Inches(1.2), Inches(3.0), Inches(1.7),
           "0.087 m", "Expert baseline\nADE", num_color=ACCENT_AMBER, num_size=Pt(36))
_stat_card(slide, Inches(10.7), Inches(1.2), Inches(1.8), Inches(1.7),
           "Ep 14", "Final ckpt\nused", num_color=MED_GRAY, num_size=Pt(28))

bullets = [
    "ADE at epoch 0 (CARLA checkpoint): 0.114 m \u2192 epoch 11: 0.085 m",
    "Policy converges to within 2 mm of expert-level accuracy",
    "Losses decrease rapidly in first 3\u20135 epochs, then plateau",
    "Language loss near zero throughout (no NL supervision)",
]
y = Inches(3.3)
for b in bullets:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(17)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    y += Inches(0.5)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 32. BASELINE ROUTE FOLLOWING ------
_table_slide(
    "Baseline Route Following (No Obstacles)",
    ["Metric", "SimLingo", "ACC Baseline"],
    [
        ["Route coverage (%)", "99.1", "97.6"],
        ["Avg lateral deviation (m)", "0.056", "0.261"],
        ["Max lateral deviation (m)", "0.311", "0.670"],
        ["Avg completion time (s)", "68.1", "42.3"],
    ],
    col_widths=[Inches(4.5), Inches(3.75), Inches(3.75)],
)


# ------ 33-37. FIGURE SLIDES ------
_figure_slide("Overall Pass Rates", FIGURES / "pass_fail_summary.png",
              "SimLingo: 60% overall (9/15), 40% obstacle-only (4/10). ACC baseline: 100%.")

_figure_slide("Route Coverage Comparison", FIGURES / "route_coverage_comparison.png",
              "Route coverage by scenario. Both controllers achieve >97% on obstacle-free runs.")

_figure_slide("Safety Comparison", FIGURES / "safety_comparison.png",
              "Collision rate, stop success rate, and stopping distance across all obstacle variants.")

_figure_slide("Lateral Deviation Comparison", FIGURES / "lateral_deviation_comparison.png",
              "SimLingo: 0.03\u20130.09 m on most scenarios. ACC: 0.25\u20130.38 m (pure-pursuit corner-cutting).")

_figure_slide("Trajectory Overlays (Bird\u2019s-Eye View)", FIGURES / "trajectory_overlays.png",
              "Bird\u2019s-eye trajectories for all variants. SimLingo (blue) vs ACC (orange) vs route (dashed).")


# ------ 38. OBSTACLE RESULTS TABLE ------
_table_slide(
    "Obstacle Avoidance: Detailed Results",
    ["Scenario", "Controller", "Coverage%", "Collisions", "Stopped", "Stop Dist (m)", "Avg Lat Dev (m)"],
    [
        ["Var 1 (Early)", "SimLingo", "23.5", "0/2", "2/2", "6.8", "0.053"],
        ["Var 1 (Early)", "ACC", "23.5", "0/2", "2/2", "7.4", "0.256"],
        ["Var 2 (Mid)", "SimLingo", "41.2", "2/2", "2/2", "3.2", "0.033"],
        ["Var 2 (Mid)", "ACC", "38.8", "0/2", "2/2", "5.4", "0.383"],
        ["Var 3 (Exit)", "SimLingo", "52.4", "0/2", "2/2", "5.7", "0.046"],
        ["Var 3 (Exit)", "ACC", "49.4", "0/2", "2/2", "8.3", "0.350"],
        ["Var 4 (Straight)", "SimLingo", "73.5", "2/2", "2/2", "3.6", "0.072"],
        ["Var 4 (Straight)", "ACC", "68.2", "0/2", "2/2", "8.1", "0.262"],
        ["Var 5 (Late)", "SimLingo", "96.5", "2/2", "0/2", "N/A", "0.284"],
        ["Var 5 (Late)", "ACC", "80.0", "0/2", "2/2", "7.3", "0.254"],
    ],
    col_widths=[Inches(2.0), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.5), Inches(1.8), Inches(1.9)],
)


# ------ 39. PER-VARIANT: CLEAN STOPS ------
_two_figure_slide(
    "Clean Stops: Variant 1 (Early) & Variant 3 (Exit)",
    FIGURES / "near_obstacle_waypoints_obstacle_var1.png",
    FIGURES / "near_obstacle_waypoints_obstacle_var3.png",
    "Var 1: Stop at ~6.8 m, speed < 1.5 m/s",
    "Var 3: Stop at ~5.7 m, just exited curve",
)


# ------ 40. PER-VARIANT: CONTACT & FAILURE ------
slide = _new_slide()
_add_title_bar(slide, "Contact & Failure: Variants 2, 4, and 5")

third_w = Inches(3.9)
third_gap = Inches(0.25)
total_figs_w = third_w * 3 + third_gap * 2
start_left = int((SLIDE_W - total_figs_w) / 2)

figs_3 = [
    (FIGURES / "pre_collision_waypoints_obstacle_var2.png", "Var 2: Low-speed contact at 3.2 m"),
    (FIGURES / "pre_collision_waypoints_obstacle_var4.png", "Var 4: Contact at 0.6 m/s, stop by step 167"),
    (FIGURES / "pre_collision_waypoints_obstacle_var5.png", "Var 5: No deceleration, 3.4 m/s through"),
]
for idx, (fig_path, cap) in enumerate(figs_3):
    region_left = int(start_left + idx * (third_w + third_gap))
    if Path(fig_path).exists():
        w, h = _fit_image(fig_path, third_w, FIG_MAX_H - Inches(0.6))
        fig_left = int(region_left + (third_w - w) / 2)
        fig_top = int(FIG_TOP + (FIG_MAX_H - Inches(0.6) - h) / 2)
        slide.shapes.add_picture(str(fig_path), fig_left, fig_top, w, h)
    cb = slide.shapes.add_textbox(region_left, SLIDE_H - Inches(0.75), third_w, Inches(0.6))
    ctf = cb.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = cap
    cp.font.size = Pt(12)
    cp.font.italic = True
    cp.font.color.rgb = MED_GRAY
    cp.alignment = PP_ALIGN.CENTER
_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ SECTION: LIMITATIONS ------
_add_section_divider("Limitations & Future Work")


# ------ 42. LIMITATIONS ------
slide = _new_slide()
_add_title_bar(slide, "Limitations")

limitations = [
    ("Sim-to-Sim Gap", "CARLA photorealistic vs QLabs simplified\u2014unclear how much visual grounding transfers", UA_RED),
    ("Timing Sensitivity", "PID derivative factor (1/5) is hand-tuned; speed from single-frame displacement is noisy", ACCENT_AMBER),
    ("Scenario Coverage", "Limited to roundabout + 5 obstacle positions; no generalization guarantee", ACCENT_BLUE),
    ("Multi-Actor", "Serialized QLabs API limits real-time multi-actor execution", MED_GRAY),
    ("No Language Supervision", "Fine-tuning adapts visual/control but not language-following", MED_GRAY),
    ("In-Distribution Eval", "Test positions overlap with training\u2014no held-out generalization test", UA_RED),
]

for i, (ltitle, ldesc, lcolor) in enumerate(limitations):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.2) + row * Inches(1.9)
    _card_with_title(slide, x, y, Inches(5.8), Inches(1.6), ltitle,
                     [ldesc], title_color=lcolor)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 43. FUTURE WORK ------
slide = _new_slide()
_add_title_bar(slide, "Future Work")

future = [
    ("Broader Scenarios", "Intersections, merges, adverse weather/lighting", ACCENT_BLUE),
    ("Multi-Actor", "Reduce lock contention, async logging", ACCENT_TEAL),
    ("Camera Calibration", "Systematic FOV/resolution matching", ACCENT_AMBER),
    ("PID Robustness", "Gain tuning sweeps across domains", MED_GRAY),
    ("Training Ablations", "LoRA rank, learning rate, augmentation", ACCENT_BLUE),
    ("Language Supervision", "Paired language annotations for commentary", ACCENT_TEAL),
    ("Sim-to-Real", "Deploy to physical QCar 2 hardware", UA_RED),
]

x = Inches(0.8)
y = Inches(1.2)
card_w = Inches(3.7)
card_h = Inches(1.35)
for i, (ftitle, fdesc, fcolor) in enumerate(future):
    col = i % 3
    row = i // 3
    cx = Inches(0.8) + col * Inches(4.1)
    cy = Inches(1.2) + row * Inches(1.65)
    _card_with_title(slide, cx, cy, Inches(3.8), Inches(1.35), ftitle, [fdesc], title_color=fcolor)

_highlight_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5),
               "Highest priority: broader scenario coverage and sim-to-real transfer to physical QCar 2",
               bg=ACCENT_RED_LIGHT, font_size=Pt(16), font_color=UA_RED, bold=True)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ SECTION: CONCLUSION ------
_add_section_divider("Conclusion")


# ------ 45. CONTRIBUTIONS ------
slide = _new_slide()
_add_title_bar(slide, "Contributions")

contribs = [
    ("1", "QLabs Expert Data Collection Pipeline",
     "Keyboard teleop, 4 Hz logging, multi-actor scene management", ACCENT_BLUE),
    ("2", "Fine-Tuning Config for InternVL2-1B with LoRA",
     "PyTorch Lightning + DeepSpeed, ~17.6M adapter parameters", ACCENT_TEAL),
    ("3", "Real-Time Inference Stack for QCar 2",
     "PID control conversion, timing compensation, ego-frame target points", ACCENT_AMBER),
    ("4", "Scenario-Driven Evaluation Harness",
     "15-run test matrix, quantitative metrics, LiDAR ACC baseline comparison", UA_RED),
]

for i, (num, ctitle, cdesc, ccolor) in enumerate(contribs):
    y = Inches(1.2) + i * Inches(1.45)
    _rounded_rect(slide, Inches(0.8), y, Inches(0.7), Inches(1.15), ccolor)
    nb = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.2), Inches(0.7), Inches(0.6))
    ntf = nb.text_frame
    np_ = ntf.paragraphs[0]
    np_.text = num
    np_.font.size = Pt(26)
    np_.font.bold = True
    np_.font.color.rgb = WHITE
    np_.font.name = FONT_HEADING
    np_.alignment = PP_ALIGN.CENTER
    # Title
    tb = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.1), Inches(10.5), Inches(0.45))
    ttf = tb.text_frame
    tp = ttf.paragraphs[0]
    tp.text = ctitle
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.color.rgb = BODY_COLOR
    tp.font.name = FONT_HEADING
    # Desc
    db = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.55), Inches(10.5), Inches(0.45))
    dtf = db.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = cdesc
    dp.font.size = Pt(15)
    dp.font.color.rgb = MED_GRAY
    dp.font.name = FONT_BODY

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 46. SUMMARY ------
slide = _new_slide()
_add_title_bar(slide, "Summary of Key Results")

# Big stat cards
_stat_card(slide, Inches(0.8), Inches(1.15), Inches(3.8), Inches(1.8),
           "25%", "ADE reduction\n0.114 \u2192 0.085 m over 15 epochs", num_color=ACCENT_BLUE, num_size=Pt(52))
_stat_card(slide, Inches(4.9), Inches(1.15), Inches(3.8), Inches(1.8),
           "60%", "Overall pass rate\n9/15 runs (5/5 baseline + 4/10 obstacle)", num_color=ACCENT_TEAL, num_size=Pt(52))
_stat_card(slide, Inches(9.0), Inches(1.15), Inches(3.8), Inches(1.8),
           "4.7\u00d7", "Lower lateral dev\nvs ACC on curves (0.056 vs 0.261 m)", num_color=ACCENT_AMBER, num_size=Pt(52))

bullets_summary = [
    "Obstacle detection and deceleration in 4 of 5 variants",
    "40% obstacle avoidance pass rate (clean stops) vs ACC baseline 100%",
    "Var 5 (high speed + curved approach): complete failure \u2014 model maintains 3\u20134 m/s",
    "Policy ADE converges to within 2 mm of expert demonstration baseline (0.087 m)",
]
y = Inches(3.3)
for b in bullets_summary:
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b
    p.font.size = Pt(17)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_BODY
    y += Inches(0.5)

_highlight_box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.55),
               "VLA domain transfer via LoRA is feasible with a small dataset \u2014 obstacle handling needs further work",
               bg=ACCENT_RED_LIGHT, font_size=Pt(17), font_color=UA_RED, bold=True)

_add_logo(slide)
_slide_number_placeholder(slide, _slide_num[0])


# ------ 47. QUESTIONS ------
slide = _new_slide()
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK

# Left accent stripe
stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_H)
stripe.fill.solid()
stripe.fill.fore_color.rgb = UA_RED
_no_line(stripe)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.25), SLIDE_W, Inches(0.04))
bar.fill.solid()
bar.fill.fore_color.rgb = UA_RED
_no_line(bar)

txBox = slide.shapes.add_textbox(Inches(0), Inches(1.8), SLIDE_W, Inches(1.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT_TITLE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Questions?"
p2.font.size = Pt(36)
p2.font.color.rgb = UA_RED
p2.font.name = FONT_HEADING
p2.font.bold = True
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(14)

info = slide.shapes.add_textbox(Inches(0), Inches(4.6), SLIDE_W, Inches(1.5))
itf = info.text_frame
for line, sz in [
    ("Garegin Mazmanyan", Pt(22)),
    ("The University of Arizona  \u00b7  Department of Computer Science", Pt(16)),
]:
    pi = itf.add_paragraph()
    pi.text = line
    pi.font.size = sz
    pi.font.color.rgb = MED_GRAY
    pi.font.name = FONT_BODY
    pi.alignment = PP_ALIGN.CENTER
    pi.space_after = Pt(4)

_add_logo(slide)


# ===================================================================
# Save
# ===================================================================
prs.save(str(OUTPUT))
print(f"Presentation saved to {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
